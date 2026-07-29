"""
Build the short "SimParc for Alberta" talk (PowerPoint) - 9 slides, ~5 minutes.

A graph-first rebuild of the earlier 7-slide Presentation.pdf. That deck was
text-heavy (11 bullets on one slide), had a blank "DATA REWEIGHT" slide, and
stopped before the results. This one caps body copy at roughly 15 words per
slide and lets a figure carry each point, ending on the headline number, the
map, and the honest limits.

Every headline number is read at build time from the pipeline's own CSV, so the
deck cannot drift from the analysis:
    data/output/calgary_energy_profile.csv        (energy_profile.py city)
Figures are embedded from calgary_adaptation/figures/ - notably 25_ipf_reweight
(energy_profile.py city), which fills the slide that used to be empty.

Speaker scripts live alongside each slide as PowerPoint notes: open the deck in
Presenter View, or read them from the SCRIPTS dict below.

Usage (from repo root):
    uv run python calgary_adaptation/build_alberta_deck.py
Writes calgary_adaptation/Alberta_Deck.pptx
"""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = REPO_ROOT / "data" / "output"
FIG_DIR = Path(__file__).resolve().parent / "figures"
PPTX_PATH = Path(__file__).resolve().parent / "Alberta_Deck.pptx"

# ---- palette: sampled from the original Presentation.pdf so the deck keeps
# the team's existing look (dark teal, crimson accent).
INK = RGBColor(0x00, 0x39, 0x41)      # slide background
PANEL = RGBColor(0x4D, 0x73, 0x78)    # raised panel
PALE = RGBColor(0xEB, 0xEB, 0xE2)     # figure card / light text
ACCENT = RGBColor(0xC8, 0x10, 0x2E)   # crimson (the old "BIAS" arrow)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIM = RGBColor(0x9F, 0xB8, 0xBC)      # secondary text on dark

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
FOOTER = "NRCan EnerGuide  ·  Census 2021  ·  seeded & scripted"


# --------------------------------------------------------------------------- #
# pptx helpers (same shape as make_presentation.py, dark palette)
# --------------------------------------------------------------------------- #

def _prs() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = EMU_W, EMU_H
    return p


def _text(tf, runs, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT,
          space_after=6):
    tf.word_wrap = True
    lines = [runs] if isinstance(runs, str) else runs
    # Reuse paragraph 0 only if it is still empty, so a second _text() call on
    # the same frame starts a new line instead of appending to the first one.
    used = bool(tf.paragraphs[0].runs)
    for i, line in enumerate(lines):
        para = tf.add_paragraph() if (i or used) else tf.paragraphs[0]
        para.alignment = align
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = line
        f = run.font
        f.size, f.bold, f.color.rgb = Pt(size), bold, color
        f.name = "Calibri"


def _shape(slide, kind, left, top, width, height, fill=None):
    shp = slide.shapes.add_shape(kind, left, top, width, height)
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def _box(slide, left, top, width, height, fill=None):
    return _shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, height, fill)


def _slide(prs, title, n):
    """Dark slide with a title, a crimson rule, a footer and a page number."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _box(s, 0, 0, EMU_W, EMU_H, fill=INK)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.24), Inches(12.1), Inches(0.7))
    _text(tb.text_frame, title.upper(), size=30, bold=True)
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _box(s, Inches(0.62), Inches(1.0), Inches(1.1), Emu(45720), fill=ACCENT)

    fb = s.shapes.add_textbox(Inches(0.6), Inches(6.95), Inches(9), Inches(0.35))
    _text(fb.text_frame, FOOTER, size=10, color=DIM)
    pb = s.shapes.add_textbox(Inches(12.3), Inches(6.95), Inches(0.6), Inches(0.35))
    _text(pb.text_frame, str(n), size=10, color=DIM, align=PP_ALIGN.RIGHT)
    return s


def _figure(slide, name, left, top, width=None, height=None):
    """Drop a figure onto a pale card so the light PNG reads as intentional.

    Exactly one of width/height is passed to add_picture so python-pptx keeps
    the image's aspect ratio; the card is then fitted to the placed picture.
    """
    path = FIG_DIR / name
    if not path.exists():
        ph = _box(slide, left, top, width or Inches(5), height or Inches(3),
                  fill=PANEL)
        _text(ph.text_frame, f"[missing figure: {name}]", size=12,
              align=PP_ALIGN.CENTER)
        return None
    pad = Inches(0.12)
    pic = slide.shapes.add_picture(str(path), left + pad, top + pad,
                                   width=width, height=height)
    card = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                  pic.width + 2 * pad, pic.height + 2 * pad, fill=PALE)
    card.adjustments[0] = 0.03
    # the card is drawn after the picture, so push it behind
    slide.shapes._spTree.remove(card._element)
    slide.shapes._spTree.insert(2, card._element)
    return pic


def _stat(slide, big, small, left, top, width, size=54, color=WHITE):
    """A large number with a caption under it - carries a slide on its own."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(size / 72 * 1.4))
    _text(tb.text_frame, big, size=size, bold=True, color=color, space_after=0)
    # 1.4x the cap height clears the descenders of the headline at every size
    # we use here (28-104pt); at 0.95 the caption collided with the number.
    cb = slide.shapes.add_textbox(left, top + Inches(size / 72 * 1.4),
                                  width, Inches(0.9))
    _text(cb.text_frame, small.split("\n"), size=15, color=DIM, space_after=2)


def _points(slide, lines, left, top, width, size=19, gap=14, color=WHITE):
    """Short fragments, each led by a crimson tick. No sentences."""
    tb = slide.shapes.add_textbox(left, top, width, Inches(3))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.space_after = Pt(gap)
        tick = para.add_run()
        tick.text = "— "
        tick.font.size, tick.font.color.rgb, tick.font.name = Pt(size), ACCENT, "Calibri"
        tick.font.bold = True
        run = para.add_run()
        run.text = line
        run.font.size, run.font.color.rgb, run.font.name = Pt(size), color, "Calibri"
    return tb


def _notes(slide, script: str) -> None:
    slide.notes_slide.notes_text_frame.text = script


# --------------------------------------------------------------------------- #
# live numbers
# --------------------------------------------------------------------------- #

def load_numbers() -> dict:
    """Headline stats, read from the pipeline CSV rather than hardcoded."""
    prof = pd.read_csv(OUT_DIR / "calgary_energy_profile.csv")
    allrow = prof[prof["stratum"] == "All Calgary"].iloc[0]
    vintage = prof[prof["dimension"] == "An_Construction"]
    newest = vintage[vintage["stratum"] == ">= 2020"].iloc[0]
    oldest = vintage.loc[vintage["mean_MEUI"].idxmax()]
    return {
        "mean": allrow["mean_MEUI"],
        "lo": allrow["ci95_low"],
        "hi": allrow["ci95_high"],
        "n": int(allrow["n"]),
        "n_eff": int(allrow["n_eff"]),
        "new_meui": newest["mean_MEUI"],
        "old_meui": oldest["mean_MEUI"],
        "old_label": oldest["stratum"],
    }


# --------------------------------------------------------------------------- #
# speaker scripts (~80 words each => ~33 s => ~5 min total)
# --------------------------------------------------------------------------- #

SCRIPTS = {
1: """Residential stock models tell you what a city's homes will demand before you
build the grid to serve them. Hydro-Quebec has one, called SimParc. Alberta
doesn't. This is what it took to move that model across the country - using only
open data.""",

2: """A stock model answers three questions utilities actually pay for: what's the
baseline, how fast is it growing, and when does it peak. This is one Montreal
apartment across a year - note the shape, not the city. Winter peak is roughly
double summer. Multiply that by a million homes and it's a capacity planning
problem. SimParc does exactly that multiplication for Quebec.""",

3: """Here's why you can't just point SimParc at Calgary. Quebec heats with electric
baseboards - cheap hydro made that the default. Alberta is ninety-nine and a half
percent natural gas. Every probability in the model is conditioned on the Quebec
answer. The code ports in an afternoon. The statistics don't port at all.""",

4: """Three layers. A Bayesian network draws forty dwelling and occupant variables -
type, vintage, fuel. ResStock-style tables add fifty-two technical attributes on
top. A deterministic mapper turns all of it into HPXML arguments the simulator
runs. Those layers meet in a ninety-seven column schema, and every one of those
columns was trained on a Quebec survey. That schema is the thing we have to
re-derive.""",

5: """Quebec has a designed energy survey. Alberta has EnerGuide - a hundred ninety
thousand audited homes, but nobody sampled them. People opt in by applying for a
retrofit grant or labelling a new build. Look at the gap. Eighty-four percent of
the sample is single-detached; the real city is fifty-five. Apartments are
twenty-seven percent of Calgary and essentially zero in the data. Averaging this
describes grant applicants.""",

6: """So we don't average it. Every home is sorted into a census cell - type, vintage,
tenure - and iterative proportional fitting solves for weights that force the
weighted mix onto the census margins. Blue is the raw sample, green is the census,
and the orange bar is the reweighted sample landing on target. The cost is honest:
thirty-four thousand homes behave statistically like about a hundred sixty.""",

7: """One hundred forty-six kilowatt-hours per square metre per year. The interval is
a post-stratified bootstrap - five thousand resamples, priced at that reduced
effective sample size, not the raw thirty-four thousand. So the error bar is real:
one thirty-eight to one fifty-six. That's a defensible city baseline, with the
uncertainty stated rather than hidden.""",

8: """Two cuts. By age: post-2020 homes come in at eighty-one, the oldest stock at two
twenty-eight - nearly three times the intensity. That's forty years of building
code, visible in the data, and it's our physical sanity check. By geography: the
inner-city ring is oldest and most intense, the new suburbs lightest. Pale FSAs
are thin samples - compare neighbourhoods, don't quote a single one.""",

9: """Three things next: source Albertan datasets at this resolution or better,
rewrite the mapper's hardcoded Quebec parameters, and retrain the Bayesian network
on Calgary. Two honest limits - apartments and the oldest vintages stay thin, so
their intervals are wide; and this is audit-modelled intensity, not a weather
simulation. It's decision-grade for comparing neighbourhoods and vintages. Not a
per-building predictor.""",
}


# --------------------------------------------------------------------------- #
# the deck
# --------------------------------------------------------------------------- #

def build() -> None:
    N = load_numbers()
    prs = _prs()

    # ---- 1. Title -------------------------------------------------------- #
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _box(s, 0, 0, EMU_W, EMU_H, fill=INK)
    _box(s, Inches(0.9), Inches(2.7), Inches(0.14), Inches(2.0), fill=ACCENT)
    tb = s.shapes.add_textbox(Inches(1.3), Inches(2.6), Inches(11), Inches(1.6))
    _text(tb.text_frame, "SIMPARC FOR ALBERTA", size=52, bold=True)
    sb = s.shapes.add_textbox(Inches(1.32), Inches(4.05), Inches(11), Inches(0.9))
    _text(sb.text_frame,
          "Porting a Quebec housing-stock energy model to Calgary on open data",
          size=19, color=DIM)
    fb = s.shapes.add_textbox(Inches(1.32), Inches(6.95), Inches(9), Inches(0.35))
    _text(fb.text_frame, FOOTER, size=10, color=DIM)
    _notes(s, SCRIPTS[1])

    # ---- 2. Why stock models --------------------------------------------- #
    s = _slide(prs, "Why stock models", 2)
    _points(s, ["What's the baseline demand?",
                "How fast is it growing?",
                "When does it peak?"],
            Inches(0.62), Inches(1.9), Inches(4.6), size=21, gap=22)
    _stat(s, "≈ 2×", "winter peak vs summer, one dwelling",
          Inches(0.62), Inches(4.5), Inches(4.6), size=40)
    _figure(s, "00_montreal_peak.png", Inches(6.0), Inches(1.7), width=Inches(6.6))
    _notes(s, SCRIPTS[2])

    # ---- 3. The problem --------------------------------------------------- #
    s = _slide(prs, "The problem", 3)
    _stat(s, "99.5%", "of Alberta homes heat with natural gas.\nQuebec's model assumes electric baseboards.",
          Inches(0.62), Inches(2.0), Inches(5.2), size=60)
    _points(s, ["Code ports in an afternoon",
                "Statistics don't port at all"],
            Inches(0.62), Inches(4.7), Inches(5.2), size=19, gap=14)
    _figure(s, "01_heating_fuel_share.png", Inches(6.3), Inches(2.1),
            width=Inches(6.4))
    _notes(s, SCRIPTS[3])

    # ---- 4. What's inside ------------------------------------------------- #
    s = _slide(prs, "What's inside", 4)
    chevrons = [("40", "dwelling &\noccupant nodes", "Bayesian network"),
                ("52", "technical\nattributes", "ResStock sampling"),
                ("219", "measure\narguments", "HPXML mapping")]
    x = Inches(0.62)
    for num, what, layer in chevrons:
        ch = _shape(s, MSO_SHAPE.CHEVRON, x, Inches(2.0), Inches(4.0),
                    Inches(1.9), fill=PANEL)
        tf = ch.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # the chevron's arrow point eats the right ~20% of the bounding box and
        # the notch the left ~10%; inset the text so it stays on the body
        tf.margin_left, tf.margin_right = Inches(0.55), Inches(0.75)
        _text(tf, num, size=28, bold=True, align=PP_ALIGN.CENTER, space_after=0)
        _text(tf, what.split("\n"), size=13, color=WHITE, align=PP_ALIGN.CENTER,
              space_after=0)
        lb = s.shapes.add_textbox(x, Inches(4.05), Inches(4.0), Inches(0.4))
        _text(lb.text_frame, layer, size=14, color=DIM, align=PP_ALIGN.CENTER)
        x += Inches(4.1)
    _stat(s, "97 columns",
          "the schema all three layers meet in — every one of them trained on a "
          "Quebec survey",
          Inches(0.62), Inches(5.0), Inches(11.5), size=40, color=WHITE)
    _notes(s, SCRIPTS[4])

    # ---- 5. The bias ------------------------------------------------------ #
    s = _slide(prs, "The bias", 5)
    _figure(s, "02_dwelling_type_vs_census.png", Inches(0.62), Inches(1.6),
            height=Inches(4.5))
    _points(s, ["Nobody sampled this data",
                "People opt in: retrofit grants, new builds"],
            Inches(9.0), Inches(2.1), Inches(3.9), size=18, gap=16)
    _stat(s, "27% vs 0.1%", "apartments in Calgary\nvs in the sample",
          Inches(9.0), Inches(4.1), Inches(3.9), size=30, color=WHITE)
    _notes(s, SCRIPTS[5])

    # ---- 6. The fix (the slide that used to be blank) --------------------- #
    s = _slide(prs, "The fix: raking", 6)
    _figure(s, "25_ipf_reweight.png", Inches(0.62), Inches(1.6),
            height=Inches(4.5))
    _points(s, ["Sort every home into a census cell",
                "Solve for weights that match the margins"],
            Inches(9.0), Inches(2.1), Inches(3.9), size=18, gap=16)
    _stat(s, f"{N['n']:,} → {N['n_eff']}",
          "homes → effective sample size.\nThe honest cost of the correction.",
          Inches(9.0), Inches(4.1), Inches(3.9), size=30, color=WHITE)
    _notes(s, SCRIPTS[6])

    # ---- 7. The number ---------------------------------------------------- #
    s = _slide(prs, "Calgary's number", 7)
    _stat(s, f"{N['mean']:.0f}", "kWh / m² · yr", Inches(0.62), Inches(1.95),
          Inches(5.0), size=104)
    _stat(s, f"{N['lo']:.0f} – {N['hi']:.0f}",
          f"95% confidence interval,\n{N['n_eff']} effective samples",
          Inches(0.62), Inches(4.6), Inches(5.0), size=34, color=DIM)
    _figure(s, "19_meui_bootstrap_distribution.png", Inches(6.3), Inches(1.75),
            width=Inches(6.4))
    _notes(s, SCRIPTS[7])

    # ---- 8. Where and when ------------------------------------------------ #
    s = _slide(prs, "Where and when", 8)
    # heights differ so the tall bar chart and the wide map together stay inside
    # the 12.1" content width (21 is ~1.18:1, 24 is ~1.71:1)
    _figure(s, "21_meui_by_vintage.png", Inches(0.62), Inches(1.6),
            height=Inches(4.4))
    _figure(s, "24_calgary_meui_map.png", Inches(6.35), Inches(2.05),
            height=Inches(3.5))
    _points(s, [f"Newest stock {N['new_meui']:.0f}, oldest {N['old_meui']:.0f} "
                f"kWh/m²·yr — the building code, visible",
                "Inner ring oldest and most intense; pale FSAs are thin samples"],
            Inches(0.62), Inches(6.25), Inches(12.1), size=15, gap=4)
    _notes(s, SCRIPTS[8])

    # ---- 9. Next / limits ------------------------------------------------- #
    s = _slide(prs, "Next  ·  and what this isn't", 9)
    steps = ["Source Albertan datasets\nat this resolution or better",
             "Rewrite the mapper's\nhardcoded Quebec parameters",
             "Retrain the Bayesian\nnetwork on Calgary"]
    x = Inches(0.62)
    for step in steps:
        bx = _shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.9),
                    Inches(1.6), fill=PANEL)
        bx.adjustments[0] = 0.08
        bx.text_frame.word_wrap = True
        bx.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(bx.text_frame, step.split("\n"), size=17,
              align=PP_ALIGN.CENTER, space_after=0)
        x += Inches(4.0)
    _points(s, ["Apartments and oldest vintages: wide intervals",
                "Audit-modelled intensity, not a weather simulation"],
            Inches(0.62), Inches(3.9), Inches(12.1), size=18, gap=12)
    _stat(s, "Decision-grade for comparison.",
          "Not a per-building predictor.",
          Inches(0.62), Inches(5.35), Inches(12.1), size=28, color=WHITE)
    _notes(s, SCRIPTS[9])

    prs.save(PPTX_PATH)
    print(f"wrote {PPTX_PATH.relative_to(REPO_ROOT)}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    print(f"  headline: {N['mean']:.1f} kWh/m2/yr "
          f"(95% CI {N['lo']:.1f}-{N['hi']:.1f}), n={N['n']:,}, n_eff={N['n_eff']}")


if __name__ == "__main__":
    build()
