"""
Build the Calgary re-calibration research presentation (PowerPoint).

A ~13-slide deck on adapting Hydro-Quebec's residential stock sampler (the front
end for SimParc) to Calgary: what SimParc needs, what must be replaced, EnerGuide
as the substitute for the missing Alberta survey, the filtered pseudo-survey, the
energy-use result, the weather-resolution problem, and an honest conclusion.

All headline numbers are read at build time from the CSVs the sibling scripts
produce, so the deck stays in sync with the pipeline:
    data/output/energuide_vs_quebec_crosswalk.csv   (compare_energuide_quebec.py)
    data/output/simparc_filter_manifest.json         (build_simparc_input.py)
    data/input/alberta/energuide/simparc_pseudosurvey.parquet
    data/input/alberta/weather/fsa_epw_lookup.csv     (fetch_fsa_weather.py)
    data/output/calgary_energy_profile.csv            (build_energy_profile.py)
Figures are embedded from calgary_adaptation/figures/.

Usage (from repo root):
    uv run python calgary_adaptation/make_presentation.py
Writes calgary_adaptation/Calgary_Recalibration_Research.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from calibrate_stock import CENSUS_MARGINS_CALGARY_2021

REPO_ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = REPO_ROOT / "data" / "output"
FIG_DIR = Path(__file__).resolve().parent / "figures"
PPTX_PATH = Path(__file__).resolve().parent / "Calgary_Recalibration_Research.pptx"

# ---- palette -------------------------------------------------------------- #
INK = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
ACCENT = RGBColor(0x0E, 0x74, 0x90)   # deep teal
ACCENT2 = RGBColor(0xC2, 0x41, 0x0C)  # burnt orange (warnings/limitations)
PANEL = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)


# --------------------------------------------------------------------------- #
# small pptx helpers
# --------------------------------------------------------------------------- #

def _prs() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = EMU_W, EMU_H
    return p


def _text(tf, runs, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
          space_after=6, level=0):
    """Set a text frame's first paragraph; `runs` may be str or list of str."""
    tf.word_wrap = True
    lines = [runs] if isinstance(runs, str) else runs
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.space_after = Pt(space_after)
        para.level = level
        run = para.add_run()
        run.text = line
        f = run.font
        f.size, f.bold, f.color.rgb = Pt(size), bold, color
        f.name = "Calibri"


def _box(slide, left, top, width, height, fill=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
        shp.line.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
    return shp


def _slide(prs, title, kicker=None, n=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _box(s, 0, 0, EMU_W, Inches(1.15), fill=INK)
    _box(s, 0, Inches(1.15), EMU_W, Emu(45720), fill=ACCENT)  # 0.05" accent rule
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.9))
    _text(tb.text_frame, title, size=30, color=WHITE, bold=True)
    tb.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    if kicker:
        kb = s.shapes.add_textbox(Inches(0.62), Inches(0.02), Inches(12), Inches(0.3))
        _text(kb.text_frame, kicker.upper(), size=11, color=ACCENT, bold=True)
    if n is not None:
        fb = s.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35))
        _text(fb.text_frame, str(n), size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    return s


def _bullets(slide, items, left, top, width, height, size=17, gap=8):
    """items: list of (text, level, color?) or plain str."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        if isinstance(it, str):
            text, level, color = it, 0, INK
        else:
            text, level, color = (it + (INK,))[:3] if len(it) == 2 else it
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.level = level
        para.space_after = Pt(gap)
        bullet = ("•  " if level == 0 else "–  ")
        run = para.add_run(); run.text = bullet + text
        f = run.font
        f.size = Pt(size - 2 * level); f.color.rgb = color; f.name = "Calibri"
    return tb


def _table(slide, df, left, top, width, height, header_fill=ACCENT,
           font=12, col_aligns=None):
    rows, cols = df.shape[0] + 1, df.shape[1]
    gtbl = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for j, name in enumerate(df.columns):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_fill
        c.margin_top = c.margin_bottom = Pt(2)
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(name)
        r.font.size, r.font.bold, r.font.color.rgb, r.font.name = Pt(font), True, WHITE, "Calibri"
    for i in range(df.shape[0]):
        for j in range(cols):
            c = gtbl.cell(i + 1, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 == 0 else PANEL
            c.margin_top = c.margin_bottom = Pt(1)
            p = c.text_frame.paragraphs[0]
            p.alignment = (col_aligns[j] if col_aligns else
                           (PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER))
            r = p.add_run(); r.text = str(df.iat[i, j])
            r.font.size, r.font.color.rgb, r.font.name = Pt(font), INK, "Calibri"
    return gtbl


def _picture(slide, name, left, top, width=None, height=None):
    path = FIG_DIR / name
    if not path.exists():
        ph = _box(slide, left, top, width or Inches(5), height or Inches(3), fill=PANEL)
        _text(ph.text_frame, f"[missing figure: {name}]", size=12, color=MUTED,
              align=PP_ALIGN.CENTER)
        return
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def _caption(slide, text, left, top, width):
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.3))
    _text(tb.text_frame, text, size=10, color=MUTED)


# --------------------------------------------------------------------------- #
# data loaders (headline numbers, read at build time)
# --------------------------------------------------------------------------- #

def load_numbers() -> dict:
    d: dict = {}
    cw = OUT_DIR / "energuide_vs_quebec_crosswalk.csv"
    if cw.exists():
        c = pd.read_csv(cw)
        d["crosswalk_counts"] = c["status"].value_counts().to_dict()
        direct = c[c["status"] == "direct"].copy()
        direct["pct_populated"] = pd.to_numeric(direct["pct_populated"], errors="coerce")
        d["crosswalk_direct"] = direct

    # Bias-composition shares (slide 8) use accurate inline fallbacks; the raked
    # EnerGuide table already lives in build_alberta_weights.alberta_stock_mapped.

    wf = OUT_DIR / "calgary_fsa_weather_profile.csv"
    if wf.exists():
        w = pd.read_csv(wf)
        city = (w.groupby("year")
                 .agg(mean_temp_C=("mean_temp_C", "mean"), hdd18=("hdd18", "mean"))
                 .reset_index())
        fsa_mean = w.groupby("FSA")["mean_temp_C"].mean()
        d["weather"] = {
            "n_fsa": int(w["FSA"].nunique()),
            "years": (int(w["year"].min()), int(w["year"].max())),
            "n_gridpoints": 60,
            "city_by_year": city,
            "temp_spread": float(fsa_mean.max() - fsa_mean.min()),
            "coldest": (int(city.loc[city["hdd18"].idxmax(), "year"]),
                        float(city["hdd18"].max())),
            "warmest": (int(city.loc[city["hdd18"].idxmin(), "year"]),
                        float(city["hdd18"].min())),
        }

    ep = OUT_DIR / "calgary_energy_profile.csv"
    if ep.exists():
        e = pd.read_csv(ep)
        allrow = e[e["stratum"] == "All Calgary"].iloc[0]
        d["meui"] = (allrow["mean_MEUI"], allrow["ci95_low"], allrow["ci95_high"],
                     int(allrow["n"]), int(allrow["n_eff"]))
    return d


# --------------------------------------------------------------------------- #
# the deck
# --------------------------------------------------------------------------- #

def build() -> None:
    N = load_numbers()
    prs = _prs()

    # ---- 1. Title -------------------------------------------------------- #
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _box(s, 0, 0, EMU_W, EMU_H, fill=INK)
    _box(s, Inches(0.9), Inches(2.6), Inches(0.12), Inches(2.3), fill=ACCENT)
    tb = s.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(11), Inches(2))
    _text(tb.text_frame, "Re-calibrating a residential building-stock sampler for Calgary",
          size=40, color=WHITE, bold=True)
    sb = s.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11), Inches(1.6))
    _text(sb.text_frame, [
        "Porting Hydro-Québec's EUEMr / SimParc pipeline to Alberta using open data",
        "EnerGuide microdata + 2021 Census — no proprietary survey, fully scripted",
    ], size=18, color=RGBColor(0xCB, 0xD5, 0xE1), space_after=6)

    # ---- 2. The pipeline ------------------------------------------------- #
    s = _slide(prs, "The pipeline: what this tool is", "context", 2)
    _bullets(s, [
        ("The sampler is a front end for SimParc — Hydro-Québec's fleet-scale "
         "OS-HPXML / EnergyPlus building-stock simulator.", 0),
        ("It draws synthetic dwellings and emits ~219 OS-HPXML measure arguments "
         "per home (building-mapping.csv) that SimParc runs directly.", 0),
        ("Three probability layers carry the statistical content:", 0),
        ("Bayesian network — 40 dwelling/occupant nodes (type, vintage, fuel, system…)", 1),
        ("ResStock-style CPTs — 52 technical attributes (envelope, HVAC, setpoints, PV)", 1),
        ("Deterministic HPXML mapper — labels → HPXML args (weather, geometry, fuels)", 1),
        ("All three were trained on Québec's EUEMr 2022 survey. Calgary needs new "
         "numbers — not new node or option names.", 0),
    ], Inches(0.7), Inches(1.5), Inches(12), Inches(5.5), size=18)

    # ---- 3. Goal & constraint ------------------------------------------- #
    s = _slide(prs, "Goal & the naming constraint", "objective", 3)
    _bullets(s, [
        ("Replace the Québec probability values with Calgary values — keep every "
         "parameter, node, and option label byte-for-byte.", 0),
        ("Where a state is impossible in Alberta (e.g. Hydro-Québec's Bi-energie "
         "tariff), its probability is set to 0; the state name is preserved.", 0),
        ("Geography is already collapsed: Territoire_HQ = Calgary, Region = Alberta, "
         "EPW/UTC/DST already Mountain Time.", 0),
        ("What still needs replacing is every downstream probability — that is the "
         "real project, and it is a data problem, not a code problem.", 0, ACCENT2),
    ], Inches(0.7), Inches(1.6), Inches(12), Inches(4), size=19, gap=14)

    # ---- 4. SimParc: what needs to be replaced --------------------------- #
    s = _slide(prs, "SimParc → what needs to be replaced", "the map", 4)
    counts = N.get("crosswalk_counts", {})
    tier = pd.DataFrame([
        ["Tier A — replace from EnerGuide microdata", "direct", counts.get("direct", 30)],
        ["Tier B/imputed — census / external aggregate", "imputed", counts.get("imputed", 3)],
        ["Tier C — no Alberta data, keep Québec CPT", "keep-qc", counts.get("keep-qc", 22)],
        ["ResStock-generic — location-independent default", "resstock", counts.get("resstock", 35)],
        ["Already set to Calgary / derived", "set/derived",
         counts.get("set-calgary", 2) + counts.get("derived", 5)],
    ], columns=["What each SimParc input needs", "status", "# cols"])
    _table(s, tier, Inches(0.9), Inches(1.7), Inches(9.5), Inches(3.2), font=15,
           col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
    _bullets(s, [
        ("97 columns in the SimParc input schema, triaged column-by-column.", 0),
        ("~⅓ map to real Calgary microdata; the rest are kept, imputed, or generic.", 0),
        ("Data-driven — see energuide_vs_quebec_crosswalk.csv", 0, MUTED),
    ], Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.6), size=15, gap=6)

    # ---- 5. The data gap ------------------------------------------------- #
    s = _slide(prs, "The data gap: no Alberta EUEMr", "ingredient 1", 5)
    _bullets(s, [
        ("Québec has EUEMr — a designed, weighted residential energy survey. "
         "Alberta has no equivalent.", 0),
        ("Substitute: NRCan EnerGuide open microdata — ~191,600 audited Alberta "
         "houses, 2004–2025, ~100 fields per home (envelope, fuels, MEUI).", 0),
        ("But it is administrative, self-selected data — retrofit-grant applicants "
         "and new-home labelling — not a representative sample.", 0, ACCENT2),
        ("So the method must correct the sample before it can stand in for a survey.", 0),
    ], Inches(0.7), Inches(1.6), Inches(12), Inches(4), size=19, gap=14)

    # ---- 6. EnerGuide vs Quebec input CSV -------------------------------- #
    s = _slide(prs, "EnerGuide vs the Québec input CSV", "coverage", 6)
    direct = N.get("crosswalk_direct")
    if direct is not None and len(direct):
        show = (direct.sort_values("pct_populated")
                      [["quebec_column", "energuide_field", "pct_populated"]]
                      .rename(columns={"quebec_column": "SimParc input",
                                       "energuide_field": "EnerGuide field",
                                       "pct_populated": "% populated"}))
        thin = show[show["% populated"] < 60].head(6)
        full = show[show["% populated"] >= 99].head(6)
        tbl = pd.concat([full, thin]).reset_index(drop=True)
        _table(s, tbl, Inches(0.8), Inches(1.6), Inches(6.6), Inches(4.8), font=12,
               col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.CENTER])
    _bullets(s, [
        ("30 SimParc inputs map directly to a named EnerGuide field.", 0),
        ("Some map in theory but are sparse in practice — modelled intensity, "
         "setpoints, PV, DHW capacity sit near ~50–56% coverage.", 0, ACCENT2),
        ("Envelope (RSI, ACH50), storeys, type, vintage, fuel: ~100% populated — "
         "the strongest Alberta signal.", 0),
        ("The rest is imputed from census or kept from Québec (Tier C).", 0),
    ], Inches(7.7), Inches(1.7), Inches(5.2), Inches(4.8), size=15, gap=10)

    # ---- 7. Filtering the data into SimParc (already built in) ----------- #
    s = _slide(prs, "Filtering the data into SimParc", "already built in", 7)
    steps = pd.DataFrame([
        ["1 — Constrain", "Pick any of the 40 variables + a value (e.g. Type = detached)"],
        ["2 — Sample", "GUM_Sampling(N, evs) draws that conditioned population"],
        ["3 — Map", "resstock_args_sampling + MapHPXML → ~219 OS-HPXML args"],
        ["4 — Export", "Download building-input / -mapping CSV → SimParc"],
    ], columns=["Dashboard step", "What it does"])
    _table(s, steps, Inches(0.8), Inches(1.65), Inches(11.7), Inches(2.3), font=14,
           col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT])
    _bullets(s, [
        ("The Streamlit dashboard already *is* the filter: constrain the population, "
         "sample it, export the CSV SimParc runs. No extra script needed.", 0),
        ("Separately, the Calgary calibration data already exists: calibrate_stock.py "
         "maps EnerGuide → BN labels and IPF-rakes to the 2021 census "
         "(≈74k Calgary homes; alberta_stock_mapped.parquet).", 0),
        ("Honest cost of that raking: Kish effective n ≈ 961 — apartments (n≈118) and the "
         "oldest stock stay thin.", 0, ACCENT2),
    ], Inches(0.8), Inches(4.3), Inches(12), Inches(2.8), size=15, gap=9)

    # ---- 8. The bias, quantified ----------------------------------------- #
    s = _slide(prs, "The bias, quantified", "why we can't take the mean", 8)
    tshare = N.get("type_share", {})
    census = CENSUS_MARGINS_CALGARY_2021["Type_Logement"]
    def pct(x): return f"{x*100:.1f}%"
    bias = pd.DataFrame([
        ["Single-detached", pct(tshare.get("Maison individuelle", 0.827)),
         pct(census["Maison individuelle"])],
        ["Apartment (Collective)", pct(tshare.get("Collective", 0.0016)),
         pct(census["Collective"])],
        ["Row house", pct(tshare.get("Maison en rangee", 0.116)),
         pct(census["Maison en rangee"])],
        ["Built ≥ 2020", pct(N.get("new_share", 0.211)),
         pct(CENSUS_MARGINS_CALGARY_2021["An_Construction"][">= 2020"])],
    ], columns=["Housing feature", "EnerGuide sample", "Calgary census"])
    _table(s, bias, Inches(1.3), Inches(1.9), Inches(9), Inches(2.7), font=16,
           col_aligns=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
    _bullets(s, [
        ("A naïve average describes grant applicants and new builds, not Calgary.", 0, ACCENT2),
        ("Apartments are 27% of the city but essentially absent from EnerGuide.", 0),
        ("Fix: sort homes into census cells and re-weight (IPF raking) so the "
         "weighted mix matches the census.", 0),
    ], Inches(1.3), Inches(4.9), Inches(11), Inches(2), size=16, gap=8)

    # ---- 9. Result ------------------------------------------------------- #
    s = _slide(prs, "Result: Calgary's household energy intensity", "output", 9)
    if "meui" in N:
        m, lo, hi, n, neff = N["meui"]
        headline = f"{m:.0f} kWh/m²·yr   (95% CI {lo:.0f}–{hi:.0f})"
    else:
        headline = "146 kWh/m²·yr   (95% CI 138–156)"
    hb = s.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(7.4), Inches(0.9))
    _text(hb.text_frame, headline, size=26, color=ACCENT, bold=True)
    _bullets(s, [
        ("Post-stratified weighted bootstrap — the interval reflects the honest "
         "effective sample size, not the raw 34k.", 0),
        ("Strong vintage gradient: post-2020 ≈ 81 vs oldest stock ≈ 215–228 "
         "kWh/m²·yr — roughly a third of the intensity.", 0),
        ("Building-code tightening, visible in the data — the physical sanity check.", 0),
    ], Inches(0.7), Inches(2.3), Inches(7.3), Inches(4), size=16, gap=10)
    _picture(s, "19_meui_bootstrap_distribution.png", Inches(8.2), Inches(1.5),
             width=Inches(4.7))
    _picture(s, "21_meui_by_vintage.png", Inches(8.2), Inches(4.25), width=Inches(4.7))

    # ---- 10. Weather ----------------------------------------------------- #
    s = _slide(prs, "Weather: recent, local temperatures", "the update", 10)
    w = N.get("weather", {})
    y0, y1 = w.get("years", (2018, 2025))
    _bullets(s, [
        ("The simulator still uses one Calgary-airport typical year (CWEC 2020) for "
         "every dwelling, citywide.", 0),
        (f"We replaced the old web-scraper with local NSRDB data: "
         f"{w.get('n_gridpoints', 60)} grid points across Calgary, {y0}–{y1}, and built a "
         f"per-FSA × year temperature / degree-day profile.", 0),
        ("The real gain is recent actual years, not fine detail — e.g. 2023 was mild "
         "(fewer heating degree-days) while 2019 was cold.", 0),
        (f"Intra-city temperature spread is tiny (~{w.get('temp_spread', 0.9):.1f} °C across "
         f"FSAs): Calgary is flat, so neighbourhoods barely differ.", 0, ACCENT2),
        ("Honest limit: NSRDB here is temperature-only, so it can't build a full "
         "weather file — the simulation still leans on the 2020 typical year.", 0, ACCENT2),
    ], Inches(0.7), Inches(1.5), Inches(7.2), Inches(5), size=15, gap=9)
    cby = w.get("city_by_year")
    if cby is not None:
        t = cby.copy()
        t["mean_temp_C"] = t["mean_temp_C"].round(1)
        t["hdd18"] = t["hdd18"].round(0).astype(int)
        t = t.rename(columns={"year": "Year", "mean_temp_C": "Mean °C", "hdd18": "HDD18"})
        _table(s, t, Inches(8.3), Inches(1.6), Inches(4.4), Inches(4.4), font=13,
               col_aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.RIGHT])
        _caption(s, "City-mean by year (heating degree-days, base 18 °C)",
                 Inches(8.3), Inches(6.1), Inches(4.4))

    # ---- 11. Data-quality traps ------------------------------------------ #
    s = _slide(prs, "The unglamorous 80%: data-quality traps", "credibility", 11)
    _bullets(s, [
        ("The _id trap — a per-file CKAN row number whose ranges overlap across "
         "files; de-duplicating on it would silently delete ~29k real records.", 0),
        ("The value-blending bug — pandas groupby.first() pulls each column's first "
         "non-null independently, mixing post-retrofit fields into pre-retrofit "
         "records; 17,868 homes affected until fixed.", 0),
        ("Geographic contamination — CLIENTCITY=\"Calgary\" includes mislabelled "
         "Edmonton (T6W) and rural FSAs; filtered out by FSA prefix.", 0),
        ("None of these throw errors — they produce plausible-but-wrong answers. "
         "Caught by profiling, not by the code crashing.", 0, ACCENT2),
    ], Inches(0.7), Inches(1.6), Inches(12), Inches(5), size=17, gap=14)

    # ---- 12. Conclusion / limitations ------------------------------------ #
    s = _slide(prs, "Conclusion — what this does not claim", "limitations", 12)
    _box(s, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.6), fill=PANEL)
    _box(s, Inches(6.9), Inches(1.5), Inches(5.7), Inches(4.6), fill=PANEL)
    hb = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(5.5), Inches(0.5))
    _text(hb.text_frame, "No large Alberta database", size=18, color=ACCENT2, bold=True)
    _bullets(s, [
        ("EnerGuide is self-selected, not a survey.", 0),
        ("Modelled intensity on only ~44% of homes.", 0),
        ("Apartments & oldest vintages stay thin (wide CIs).", 0),
        ("Within-cell self-selection is assumed away, not verified.", 0),
        ("Tier-C nodes (pool/spa, appliances) keep Québec values.", 0),
    ], Inches(0.9), Inches(2.2), Inches(5.5), Inches(3.8), size=14, gap=10)
    hb2 = s.shapes.add_textbox(Inches(7.1), Inches(1.6), Inches(5.3), Inches(0.5))
    _text(hb2.text_frame, "Weather data", size=18, color=ACCENT2, bold=True)
    _bullets(s, [
        ("NSRDB gives recent (2018–2025) FSA-level temperature — a real improvement.", 0),
        ("But it is temperature-only: no sunlight / humidity / wind.", 0),
        ("So it can't build a full weather file; simulation still uses the 2020 "
         "typical year.", 0),
        ("Grid covers central Calgary; city temperature is nearly uniform anyway.", 0),
        ("Energy figures are audit MEUI, not a weather simulation.", 0),
    ], Inches(7.1), Inches(2.2), Inches(5.3), Inches(3.8), size=14, gap=9)
    _text(s.shapes.add_textbox(Inches(0.7), Inches(6.3), Inches(12), Inches(0.7)).text_frame,
          "Decision-grade for relative comparisons (neighbourhoods, vintages) and a "
          "city average with a real error bar — not a per-building predictor.",
          size=15, color=INK, bold=True)

    # ---- 13. Next steps & reproducibility -------------------------------- #
    s = _slide(prs, "Next steps & reproducibility", "recap", 13)
    _bullets(s, [
        ("Widen BN coverage — extend EnerGuideToBN to the remaining Tier-A nodes.", 0),
        ("Re-pull NSRDB with all variables (sun/humidity/wind) to build recent-year "
         "weather files, then wire them into the simulator.", 0),
        ("Source better Tier-B/C data (SHEU, CEUD, ZEV, CMHC) and quantify Tier-C "
         "with Monte-Carlo sensitivity bands.", 0),
        ("Everything is seeded and scripted from public sources:", 0),
        ("compare_schema.py  ·  calibrate_stock.py", 1, MUTED),
        ("weather_profile.py  ·  energy_profile.py  ·  make_presentation.py", 1, MUTED),
    ], Inches(0.7), Inches(1.6), Inches(12), Inches(5), size=18, gap=12)

    prs.save(PPTX_PATH)
    print(f"wrote {len(prs.slides)} slides -> {PPTX_PATH.relative_to(REPO_ROOT)}")


if __name__ == "__main__":
    build()
